from pathlib import Path
from html import escape
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


ROOT = Path("/Users/dds/data/ai/ai-study/iphone18-fruit-chain-ppt")
SLIDES = ROOT / "slides"
OUT = ROOT / "output"


stocks = [
    {"code": "601138", "name": "工业富联", "seg": "组装代工", "price": 66.01, "chg": 3.50, "ret": 150, "m20": -11.7, "pos": 69, "trend": "整理", "trigger": "收复 68.55/70.93", "risk": "前期涨幅大，订单兑现敏感"},
    {"code": "002475", "name": "立讯精密", "seg": "组装/连接器", "price": 62.47, "chg": -1.28, "ret": 77, "m20": -9.9, "pos": 59, "trend": "整理", "trigger": "收复 66.88/68.55", "risk": "Pro 预期拥挤，跌破均线后修复未确认"},
    {"code": "002456", "name": "欧菲光", "seg": "光学/摄像", "price": 8.78, "chg": -3.41, "ret": -24, "m20": -8.0, "pos": 6, "trend": "弱势", "trigger": "先站回 9.14", "risk": "低位不等于反转，需量价修复"},
    {"code": "603501", "name": "韦尔股份", "seg": "CMOS/半导体", "price": 100.95, "chg": 1.50, "ret": -19, "m20": 8.2, "pos": 22, "trend": "修复", "trigger": "守 92.23，上看 105.72", "risk": "仍处年内低位区，反弹持续性待验"},
    {"code": "002273", "name": "水晶光电", "seg": "光学棱镜", "price": 32.04, "chg": -6.23, "ret": 60, "m20": -10.4, "pos": 51, "trend": "整理", "trigger": "收复 35.25/35.39", "risk": "当日大跌，短线筹码松动"},
    {"code": "300433", "name": "蓝思科技", "seg": "玻璃结构件", "price": 46.05, "chg": -11.19, "ret": 97, "m20": 2.8, "pos": 66, "trend": "回踩", "trigger": "观察 39.22 支撑", "risk": "单日大跌，需防高位补跌"},
    {"code": "002938", "name": "鹏鼎控股", "seg": "FPC/PCB", "price": 91.29, "chg": -1.23, "ret": 132, "m20": -21.0, "pos": 63, "trend": "整理", "trigger": "站回 92.37 转稳", "risk": "20 日动量明显转弱"},
    {"code": "002384", "name": "东山精密", "seg": "FPC/精密电子", "price": 237.56, "chg": 0.29, "ret": 447, "m20": 5.9, "pos": 83, "trend": "高位强势", "trigger": "244.47 短压，213.55 风险线", "risk": "涨幅极大，估值与业绩兑现压力高"},
    {"code": "002241", "name": "歌尔股份", "seg": "声学组件", "price": 20.41, "chg": -3.18, "ret": -10, "m20": -14.4, "pos": 0, "trend": "弱势", "trigger": "先收复 22.64", "risk": "年内低位弱势，不宜把低位等同安全"},
    {"code": "600584", "name": "长电科技", "seg": "封测", "price": 94.11, "chg": -6.81, "ret": 181, "m20": 25.0, "pos": 78, "trend": "强势", "trigger": "守 89.93；破位看 69.54", "risk": "强势票单日大跌，波动放大"},
    {"code": "688008", "name": "澜起科技", "seg": "芯片配套", "price": 247.15, "chg": -2.39, "ret": 198, "m20": 4.5, "pos": 66, "trend": "回踩", "trigger": "230.87 为趋势防线", "risk": "高位回踩，需防估值切换"},
    {"code": "300207", "name": "欣旺达", "seg": "电池", "price": 17.24, "chg": -2.82, "ret": -12, "m20": -14.3, "pos": 0, "trend": "弱势", "trigger": "先收复 19.42", "risk": "低位弱势，缺少趋势确认"},
    {"code": "300136", "name": "信维通信", "seg": "射频/天线", "price": 89.88, "chg": -3.77, "ret": 303, "m20": -12.7, "pos": 63, "trend": "整理", "trigger": "收复 97.90/103.25", "risk": "涨幅大且跌破均线，需防回撤扩大"},
    {"code": "002600", "name": "领益智造", "seg": "结构件辅料", "price": 15.06, "chg": -5.46, "ret": 68, "m20": 1.3, "pos": 63, "trend": "整理", "trigger": "站回 15.35/16.15", "risk": "当日大跌，资金退潮信号需观察"},
    {"code": "300709", "name": "精研科技", "seg": "散热结构件", "price": 45.95, "chg": -4.85, "ret": 9, "m20": -11.5, "pos": 30, "trend": "整理", "trigger": "先修复 48.37", "risk": "偏离 20 日线较大，波动高"},
    {"code": "603283", "name": "赛腾股份", "seg": "检测设备", "price": 67.80, "chg": -0.66, "ret": 103, "m20": 42.2, "pos": 63, "trend": "强势", "trigger": "守 65.95，破位看 59.23", "risk": "短期涨速快，回撤空间随之放大"},
]

slides = [
    ("封面", "iPhone18 预期果链观察", "A 股核心环节、标的信号与风险提示", "只做行业板块与标的信息梳理，不构成买卖建议。数据时点：2026-07-08 16:25。"),
    ("市场", "市场不是全面顺风", "半导体与 IT 设备尚有支撑，元器件/机械/电气设备承压", "上证 -0.48%，深成指 -1.87%，创业板 -1.70%；元器件 -2.45%，专用机械 -2.80%，电气设备 -3.77%。"),
    ("地图", "果链拆成五条观察线", "组装、光学、FPC、芯片封测、结构设备", "这不是单一板块行情，而是多个环节共同受 Apple 订单预期、消费电子景气和估值切换影响。"),
    ("分层", "技术状态出现明显分化", "强势票少，整理票多，弱势票仍需先修复均线", "强势：长电科技、赛腾股份；修复：韦尔股份；高位回踩：东山精密、澜起科技、蓝思科技；低位弱势：欧菲光、歌尔股份、欣旺达。"),
    ("强势", "强势票看防线，不追叙事", "高涨幅标的的核心不是故事，而是 20/60 日线是否守住", "东山精密、信维通信、澜起科技、长电科技等近一年涨幅巨大，一旦跌破趋势线，回撤弹性会显著放大。"),
    ("观察", "把买卖问题转成触发条件", "右侧确认：站回 20/60 日线；风险控制：跌破趋势防线", "不输出买入建议。只给观察条件：站回关键均线、成交配合、板块同步修复；反向条件是放量跌破 60 日线或预期兑现不及。"),
    ("风险", "风险清单比故事更重要", "订单、良率、砍单、汇率、估值、资金拥挤度", "供应链身份需以公司公告/客户认证为准；iPhone18 仍为市场预期，不应写成 Apple 官方确认事项。"),
    ("结论", "现阶段更适合做观察池", "分层跟踪，等待均线与板块共振确认", "强势票看趋势防线，弱势票看修复确认；任何个股都不因“果链”标签自动具备安全边际。"),
]

risks = [
    "Apple 官方节奏未确认：iPhone18 仍是市场预期",
    "订单与砍单风险：供应链传导快于基本面披露",
    "良率/ASP/汇率：影响单机价值量和毛利率",
    "涨幅拥挤：高位趋势票对负面信息更敏感",
    "供应链身份：以公告与客户认证为准，少用传闻定价",
]


palette = {
    "paper": "F7F4EE",
    "ink": "1D2528",
    "muted": "667074",
    "line": "D8D0C5",
    "accent": "B9532B",
    "green": "2C7A55",
    "red": "B33A3A",
    "blue": "315E72",
    "dark": "1E2A2F",
}


def css():
    return f"""
* {{ box-sizing: border-box; }}
html, body {{ margin: 0; padding: 0; }}
body {{
  width: 960pt;
  height: 540pt;
  overflow: hidden;
  background: #{palette['paper']};
  color: #{palette['ink']};
  font-family: "PingFang SC", "Microsoft YaHei", "Noto Sans CJK SC", Arial, sans-serif;
  position: relative;
}}
.mast {{ position: absolute; left: 42pt; top: 26pt; right: 42pt; height: 30pt; border-bottom: 1pt solid #{palette['line']}; }}
.mast p {{ font-size: 8.5pt; color: #{palette['muted']}; letter-spacing: .08em; margin: 0; text-transform: uppercase; }}
.page {{ position: absolute; right: 42pt; bottom: 24pt; }}
.page p {{ font-size: 8pt; color: #{palette['muted']}; margin: 0; }}
.kicker {{ position:absolute; left: 42pt; top: 86pt; width: 150pt; }}
.kicker p {{ margin: 0; font-size: 10pt; color: #{palette['accent']}; font-weight: 700; letter-spacing: .12em; }}
.title {{ position:absolute; left: 42pt; top: 120pt; width: 520pt; }}
.title h1 {{ margin:0; font-size: 36pt; line-height: 1.06; color: #{palette['ink']}; font-weight: 800; letter-spacing: 0; }}
.subtitle {{ position:absolute; left: 42pt; top: 214pt; width: 500pt; }}
.subtitle p {{ margin:0; font-size: 17pt; line-height: 1.42; color: #{palette['muted']}; }}
.note {{ position:absolute; left: 42pt; bottom: 42pt; width: 600pt; }}
.note p {{ font-size: 10.5pt; color: #{palette['muted']}; line-height: 1.45; margin: 0; }}
.panel {{ position:absolute; background: #FFFFFF; border: 1pt solid #{palette['line']}; border-radius: 4pt; padding: 14pt; }}
.panel h2 {{ margin: 0 0 7pt 0; font-size: 15pt; line-height:1.2; color: #{palette['ink']}; }}
.panel p {{ margin: 0; font-size: 10.5pt; line-height: 1.45; color: #{palette['muted']}; }}
.big-num h2 {{ font-size: 30pt; line-height: 1; margin: 0; color: #{palette['accent']}; }}
.tag {{ position:absolute; background: #{palette['dark']}; border-radius: 999pt; padding: 5pt 10pt; }}
.tag p {{ margin:0; color:white; font-size: 9pt; }}
.small {{ font-size: 9pt !important; }}
.grid-table {{ position:absolute; left:42pt; top:126pt; width:876pt; }}
.row {{ display:flex; border-bottom:1pt solid #{palette['line']}; min-height:25pt; }}
.cell {{ padding:5pt 6pt; }}
.cell p {{ margin:0; font-size:8.7pt; line-height:1.25; color:#{palette['ink']}; }}
.head {{ background: #{palette['dark']}; }}
.head p {{ color:white; font-weight:700; }}
"""


def slide_shell(idx, kind, title, subtitle, body_html):
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>{idx:02d} · {escape(kind)}</title>
<style>{css()}</style>
</head>
<body>
<!-- assumptions: 投研简报场景；使用用户上轮行情数据；不输出买入建议；iPhone18 只作为市场预期。 -->
<div class="mast"><p>IPHONE18 EXPECTATION SUPPLY CHAIN · A SHARE WATCHLIST · 2026-07-08</p></div>
<div class="kicker"><p>{escape(kind)}</p></div>
<div class="title"><h1>{escape(title)}</h1></div>
<div class="subtitle"><p>{escape(subtitle)}</p></div>
{body_html}
<div class="page"><p>{idx:02d} / 08</p></div>
</body>
</html>
"""


def metric_panel(x, y, label, value, detail, color=None):
    c = color or palette["accent"]
    return f"""<div class="panel big-num" style="left:{x}pt;top:{y}pt;width:160pt;height:92pt;">
  <h2 style="color:#{c};">{escape(value)}</h2>
  <p><b>{escape(label)}</b><br>{escape(detail)}</p>
</div>"""


def write_html_deck():
    SLIDES.mkdir(parents=True, exist_ok=True)
    OUT.mkdir(parents=True, exist_ok=True)
    (ROOT / "product-facts.md").write_text(
        "# Product Facts\n\n"
        "- Apple 官网当前可见产品线为 iPhone 17 / iPhone Air / iPhone 17e 等，并未把 iPhone 18 作为官方已发布产品展示。\n"
        "- 第三方资料对 iPhone 18 的发布时间与产品分层仍使用 rumor / reports 等表述；本 deck 只使用「iPhone18 预期」措辞。\n"
        "- 本 deck 的股票行情和技术信号来自 a-share-skill 脚本在 2026-07-08 16:25 左右返回的数据。\n",
        encoding="utf-8",
    )
    index = """<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8"><title>iPhone18 预期果链观察</title>
<style>*{box-sizing:border-box}html,body{height:100%;margin:0;background:#111;overflow:hidden;font-family:-apple-system,"PingFang SC",sans-serif}#stage{position:fixed;top:0;left:0;width:1280px;height:720px;transform-origin:top left;background:#fff;box-shadow:0 10px 60px rgba(0,0,0,.4)}iframe{width:100%;height:100%;border:0}.counter{position:fixed;right:18px;bottom:18px;color:#fff;background:rgba(0,0,0,.58);border-radius:999px;padding:6px 12px;font-size:13px}</style></head><body>
<div id="stage"><iframe id="frame"></iframe></div><div class="counter" id="counter"></div>
<script>
const deck=[
""" + ",\n".join([f'  {{file:"slides/{i:02d}-{name}.html",label:"{label}"}}' for i, (name, label) in enumerate([
        ("cover", "封面"), ("market", "市场"), ("map", "产业链地图"), ("segments", "状态分层"),
        ("leaders", "强势与回踩"), ("triggers", "观察条件"), ("risks", "风险"), ("close", "结论")
    ], 1)]) + """
];let i=0;const f=document.getElementById('frame'),c=document.getElementById('counter'),s=document.getElementById('stage');
function fit(){const z=Math.min(innerWidth/1280,innerHeight/720);s.style.transform=`translate(${(innerWidth-1280*z)/2}px,${(innerHeight-720*z)/2}px) scale(${z})`;}
function show(n){i=Math.max(0,Math.min(deck.length-1,n));f.src=deck[i].file;c.textContent=`${i+1} / ${deck.length} · ${deck[i].label}`;location.hash='#'+(i+1);}
addEventListener('resize',fit);addEventListener('keydown',e=>{if(['ArrowRight',' ','PageDown'].includes(e.key))show(i+1);if(['ArrowLeft','PageUp'].includes(e.key))show(i-1);if(e.key==='Home')show(0);if(e.key==='End')show(deck.length-1);});
fit();show((parseInt(location.hash.slice(1),10)||1)-1);
</script></body></html>"""
    (ROOT / "index.html").write_text(index, encoding="utf-8")

    bodies = []
    bodies.append(metric_panel(598, 118, "实时样本", "16/16", "实时行情成功返回") + metric_panel(760, 118, "历史样本", "243", "近一年交易日") + metric_panel(598, 222, "交付边界", "0", "不做买入建议", palette["blue"]) + metric_panel(760, 222, "结论形态", "观察池", "等待共振确认", palette["green"]) + f'<div class="note"><p>{escape(slides[0][3])}</p></div>')
    bodies.append(metric_panel(584, 128, "半导体", "+0.47%", "相对抗跌", palette["green"]) + metric_panel(750, 128, "IT设备", "+1.89%", "资金仍有承接", palette["green"]) + metric_panel(584, 236, "元器件", "-2.45%", "果链核心承压", palette["red"]) + metric_panel(750, 236, "专用机械", "-2.80%", "设备端回撤", palette["red"]) + f'<div class="note"><p>{escape(slides[1][3])}</p></div>')
    map_html = '<div style="position:absolute;left:42pt;top:302pt;width:876pt;height:150pt;">'
    segs = ["组装代工", "光学/摄像", "FPC/PCB", "芯片/封测", "结构/设备"]
    for n, seg in enumerate(segs):
        map_html += f'<div class="panel" style="left:{n*174}pt;top:0;width:158pt;height:112pt;"><h2>{seg}</h2><p>'
        names = "、".join([s["name"] for s in stocks if seg.split("/")[0] in s["seg"] or (seg=="结构/设备" and any(k in s["seg"] for k in ["结构","设备","玻璃"]))][:4])
        map_html += escape(names) + '</p></div>'
    map_html += '</div><div class="note"><p>' + escape(slides[2][3]) + '</p></div>'
    bodies.append(map_html)
    layer = '<div style="position:absolute;left:586pt;top:112pt;width:330pt;height:330pt;">'
    groups = [("强势", ["长电科技", "赛腾股份"], palette["green"]), ("修复", ["韦尔股份"], palette["blue"]), ("高位回踩", ["东山精密", "澜起科技", "蓝思科技"], palette["accent"]), ("低位弱势", ["欧菲光", "歌尔股份", "欣旺达"], palette["red"])]
    for j,(g,ns,col) in enumerate(groups):
        layer += f'<div class="panel" style="left:0;top:{j*78}pt;width:330pt;height:64pt;"><h2 style="color:#{col};">{g}</h2><p>{escape("、".join(ns))}</p></div>'
    layer += '</div><div class="note"><p>' + escape(slides[3][3]) + '</p></div>'
    bodies.append(layer)
    top = sorted(stocks, key=lambda x: x["ret"], reverse=True)[:6]
    table = '<div class="grid-table" style="top:296pt;"><div class="row head"><div class="cell" style="width:112pt;"><p>标的</p></div><div class="cell" style="width:92pt;"><p>近1年</p></div><div class="cell" style="width:110pt;"><p>状态</p></div><div class="cell" style="width:340pt;"><p>观察条件</p></div><div class="cell" style="width:222pt;"><p>主要风险</p></div></div>'
    for s in top:
        color = palette["green"] if s["ret"] > 150 else palette["accent"]
        table += f'<div class="row"><div class="cell" style="width:112pt;"><p>{s["name"]}</p></div><div class="cell" style="width:92pt;"><p style="color:#{color};font-weight:700;">+{s["ret"]}%</p></div><div class="cell" style="width:110pt;"><p>{s["trend"]}</p></div><div class="cell" style="width:340pt;"><p>{escape(s["trigger"])}</p></div><div class="cell" style="width:222pt;"><p>{escape(s["risk"])}</p></div></div>'
    table += '</div>'
    bodies.append(table)
    trig = '<div style="position:absolute;left:578pt;top:112pt;width:340pt;">'
    for j,(h,p,col) in enumerate([("右侧确认", "收盘站回 20/60 日线，且板块同步修复", palette["green"]), ("风险触发", "放量跌破 60 日线，或强势票单日大跌后不能修复", palette["red"]), ("左侧克制", "低位票先看 20 日线，不把便宜当成安全", palette["blue"])]):
        trig += f'<div class="panel" style="left:0;top:{j*96}pt;width:340pt;height:78pt;"><h2 style="color:#{col};">{h}</h2><p>{p}</p></div>'
    trig += '</div><div class="note"><p>' + escape(slides[5][3]) + '</p></div>'
    bodies.append(trig)
    rhtml = '<div style="position:absolute;left:586pt;top:112pt;width:330pt;">'
    for j, r in enumerate(risks):
        rhtml += f'<div class="panel" style="left:0;top:{j*61}pt;width:330pt;height:48pt;"><p>{escape(r)}</p></div>'
    rhtml += '</div><div class="note"><p>' + escape(slides[6][3]) + '</p></div>'
    bodies.append(rhtml)
    bodies.append(metric_panel(585, 130, "强势票", "守线", "20/60 日趋势防线", palette["green"]) + metric_panel(760, 130, "弱势票", "修复", "先收复 20 日线", palette["blue"]) + metric_panel(585, 245, "交易表达", "条件", "不输出买卖建议", palette["accent"]) + metric_panel(760, 245, "后续跟踪", "共振", "板块与个股同步", palette["green"]) + f'<div class="note"><p>{escape(slides[7][3])}</p></div>')

    for idx, ((kind, title, subtitle, _), body) in enumerate(zip(slides, bodies), 1):
        filename = ["cover", "market", "map", "segments", "leaders", "triggers", "risks", "close"][idx - 1]
        (SLIDES / f"{idx:02d}-{filename}.html").write_text(slide_shell(idx, kind, title, subtitle, body), encoding="utf-8")


def add_text(slide, x, y, w, h, text, size=18, bold=False, color="1D2528", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def add_card(slide, x, y, w, h, title, body, color="B9532B"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shp.line.color.rgb = RGBColor.from_string(palette["line"])
    add_text(slide, x + 0.12, y + 0.1, w - 0.24, 0.28, title, 14, True, color)
    add_text(slide, x + 0.12, y + 0.46, w - 0.24, h - 0.52, body, 9.6, False, palette["muted"])


def setup_slide(prs, idx, kind, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor.from_string(palette["paper"])
    add_text(slide, 0.58, 0.28, 11.9, 0.25, "IPHONE18 EXPECTATION SUPPLY CHAIN · A SHARE WATCHLIST · 2026-07-08", 8, False, palette["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.58), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(palette["line"])
    line.line.fill.background()
    add_text(slide, 0.58, 1.0, 2.2, 0.3, kind, 9, True, palette["accent"])
    add_text(slide, 0.58, 1.38, 7.15, 1.05, title, 31, True, palette["ink"])
    add_text(slide, 0.58, 2.52, 6.9, 0.75, subtitle, 14.5, False, palette["muted"])
    add_text(slide, 12.2, 7.05, 0.55, 0.2, f"{idx:02d}/08", 8, False, palette["muted"], PP_ALIGN.RIGHT)
    return slide


def write_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, (kind, title, subtitle, note) in enumerate(slides, 1):
        slide = setup_slide(prs, idx, kind, title, subtitle)
        if idx == 1:
            add_card(slide, 8.15, 1.28, 1.7, 0.9, "16/16", "实时行情成功返回")
            add_card(slide, 10.05, 1.28, 1.7, 0.9, "243", "近一年交易日")
            add_card(slide, 8.15, 2.45, 1.7, 0.9, "0", "不做买入建议", palette["blue"])
            add_card(slide, 10.05, 2.45, 1.7, 0.9, "观察池", "等待共振确认", palette["green"])
        elif idx == 2:
            for x, y, t, b, c in [(7.9, 1.55, "半导体 +0.47%", "相对抗跌", palette["green"]), (9.95, 1.55, "IT设备 +1.89%", "资金承接", palette["green"]), (7.9, 2.7, "元器件 -2.45%", "果链核心承压", palette["red"]), (9.95, 2.7, "专用机械 -2.80%", "设备端回撤", palette["red"])]:
                add_card(slide, x, y, 1.8, 0.92, t, b, c)
        elif idx == 3:
            for n, seg in enumerate(["组装代工", "光学/摄像", "FPC/PCB", "芯片/封测", "结构/设备"]):
                add_card(slide, 0.72 + n * 2.35, 4.2, 2.05, 1.18, seg, "、".join([s["name"] for s in stocks if seg.split("/")[0] in s["seg"] or (seg == "结构/设备" and any(k in s["seg"] for k in ["结构", "设备", "玻璃"]))][:4]), palette["blue"])
        elif idx == 4:
            for n, (g, ns, c) in enumerate([("强势", "长电科技、赛腾股份", palette["green"]), ("修复", "韦尔股份", palette["blue"]), ("高位回踩", "东山精密、澜起科技、蓝思科技", palette["accent"]), ("低位弱势", "欧菲光、歌尔股份、欣旺达", palette["red"])]):
                add_card(slide, 8.0, 1.25 + n * 0.95, 3.6, 0.7, g, ns, c)
        elif idx == 5:
            for n, s in enumerate(sorted(stocks, key=lambda x: x["ret"], reverse=True)[:6]):
                add_card(slide, 0.6 + (n % 3) * 4.1, 4.0 + (n // 3) * 1.05, 3.75, 0.78, f'{s["name"]} +{s["ret"]}%', f'{s["trend"]} · {s["trigger"]}', palette["green"] if s["ret"] > 150 else palette["accent"])
        elif idx == 6:
            for n, (h, b, c) in enumerate([("右侧确认", "收盘站回 20/60 日线，且板块同步修复", palette["green"]), ("风险触发", "放量跌破 60 日线，或强势票单日大跌后不能修复", palette["red"]), ("左侧克制", "低位票先看 20 日线，不把便宜当成安全", palette["blue"])]):
                add_card(slide, 8.0, 1.35 + n * 1.15, 3.9, 0.88, h, b, c)
        elif idx == 7:
            for n, r in enumerate(["Apple 官方节奏未确认", "订单与砍单风险", "良率/ASP/汇率", "涨幅拥挤", "供应链身份需验证"]):
                add_card(slide, 8.0, 1.15 + n * 0.72, 3.8, 0.52, r, risks[n], palette["red"] if n in [0, 3] else palette["accent"])
        elif idx == 8:
            for x, y, t, b, c in [(8.0, 1.4, "强势票", "看 20/60 日趋势防线", palette["green"]), (10.0, 1.4, "弱势票", "先看 20 日线修复", palette["blue"]), (8.0, 2.55, "表达", "条件，不是建议", palette["accent"]), (10.0, 2.55, "跟踪", "板块与个股共振", palette["green"])]:
                add_card(slide, x, y, 1.75, 0.9, t, b, c)
        add_text(slide, 0.58, 6.62, 8.2, 0.32, note, 9.5, False, palette["muted"])
    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(OUT / "iPhone18_A股果链观察_20260708.pptx")


if __name__ == "__main__":
    write_html_deck()
    write_pptx()
    print(ROOT)
